#!/usr/bin/env python3
"""
validate_pptx.py — Validación automática de presentaciones educativas.

Uso:
    python validate_pptx.py <path.pptx> [--grade 5to|secundaria|colegio]

Checks:
  ✓ Bloom objectives en Slide 1
  ✓ Max palabras por slide según grado
  ✓ Sin párrafos continuos (>3 líneas)
  ✓ Jerga técnica detectada (necesita glosario)
  ✓ Preguntas de reflexión presentes
  ✓ Footer con "Grado - Materia"
  ✓ Paleta correcta para el nivel
  ✓ Glosario en última slide
  ✓ Imágenes con atribución (CC)
  ✓ Progressive disclosure (fade transitions)
"""

import sys
import os
import re
from pptx import Presentation
from pptx.dml.color import RGBColor

# ── CONFIG ───────────────────────────────────────────────────────────────────
GRADE_CONFIG = {
    '5to':     {'max_words': 80,  'template': 'WARM',         'min_questions': 2},
    'secundaria': {'max_words': 120, 'template': 'PROFESSIONAL', 'min_questions': 1},
    'colegio':    {'max_words': 200, 'template': 'ACADEMIC',     'min_questions': 1},
}

TEMPLATE_PALETTES = {
    'WARM': {
        'primary': (0x0F, 0x7C, 0x7C), 'secondary': (0xE8, 0x7A, 0x2A),
        'accent': (0xF4, 0xC5, 0x42), 'bg': (0xFF, 0xF8, 0xF0),
    },
    'PROFESSIONAL': {
        'primary': (0x1B, 0x3A, 0x5C), 'secondary': (0x2E, 0x86, 0xAB),
        'accent': (0xF4, 0xA2, 0x61), 'bg': (0xF0, 0xF4, 0xF8),
    },
    'ACADEMIC': {
        'primary': (0x1A, 0x1A, 0x2E), 'secondary': (0x7C, 0x3A, 0xED),
        'accent': (0xF5, 0x9E, 0x0B), 'bg': (0xF8, 0xFA, 0xFC),
    },
}

TECH_JARGON = [
    'megafauna', 'restos liticos', 'paleolagos', 'glaciacion',
    'cultura clovis', 'tupiguarani', 'pleistoceno', 'estrato',
    'sedimento', 'artefacto litico', 'ceramica', 'osteologico',
    'bifacial', 'acanalado', 'paleoindio', 'holoceno',
]

BLOOM_LEVELS = ['recordar', 'comprender', 'aplicar', 'analizar', 'evaluar', 'crear']

TRANSITION_COLORS = {
    'TEAL': (0x0F, 0x7C, 0x7C),
    'TEAL-ish': (0x0F, 0x7C, 0x7C),
    'DK_BLUE-ish': (0x1B, 0x3A, 0x5C),
    'NAVY-ish': (0x1A, 0x1A, 0x2E),
}

# ── HELPERS ──────────────────────────────────────────────────────────────────
def rgb_dist(a, b):
    """Distancia euclidiana entre dos colores RGB."""
    return sum((x - y) ** 2 for x, y in zip(a, b)) ** 0.5

def closest_palette(slide):
    """Estima a qué template se parece más el fondo del slide."""
    bg = slide.background.fill
    if bg.type is None:
        return None
    try:
        c = bg.fore_color.rgb
        slide_rgb = (c[0], c[1], c[2])
    except:
        return None
    
    best = None
    best_dist = 999
    for name, colors in TEMPLATE_PALETTES.items():
        d = rgb_dist(slide_rgb, colors['bg'])
        if d < best_dist:
            best_dist = d
            best = name
    return best if best_dist < 100 else None

def extract_text_from_slide(slide):
    """Extrae TODO el texto de un slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return texts

def count_words(texts):
    return sum(len(t.split()) for t in texts)

def has_jargon(texts):
    """Detecta jargon técnico en el texto del slide."""
    found = []
    full_text = ' '.join(texts).lower()
    for term in TECH_JARGON:
        if term in full_text:
            found.append(term)
    return found

def has_questions(texts):
    """Detecta preguntas de reflexión."""
    q_markers = ['?', 'piensa', 'imagina', 'y tu', 'que crees', 'por que', 'como']
    full = ' '.join(texts).lower()
    return sum(1 for m in q_markers if m in full)

def has_glossary(texts):
    """Detecta glosario: línea con '=' o '=' o definición corta."""
    full = ' '.join(texts)
    indicators = ['glosario', 'terminos', '= ', 'significa', 'definicion']
    return any(i in full.lower() for i in indicators)

def has_attribution(texts):
    """Detecta atribución CC."""
    full = ' '.join(texts).lower()
    markers = ['wikimedia', 'cc by', 'creative commons', 'attribution', 'fuente:']
    return any(m in full for m in markers)

def has_bloom_levels(texts):
    """Detecta objetivos Bloom."""
    full = ' '.join(texts).lower()
    found = [l for l in BLOOM_LEVELS if l in full]
    return found

def has_footer(texts):
    """Detecta footer con grado+materia."""
    full = ' '.join(texts)
    return bool(re.search(r'\dto\s.*(primaria|secundaria|colegio)', full, re.I))

def has_fade(slide):
    """Detecta si el slide tiene transición fade."""
    from lxml import etree
    ns = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    trans = slide._element.find(f'{ns}transition')
    if trans is not None:
        fade = trans.find(f'{ns}fade')
        return fade is not None
    return False

# ── MAIN VALIDATOR ──────────────────────────────────────────────────────────
def validate(pptx_path, grade='5to'):
    if not os.path.exists(pptx_path):
        return {'error': f'File not found: {pptx_path}'}
    
    config = GRADE_CONFIG.get(grade, GRADE_CONFIG['5to'])
    prs = Presentation(pptx_path)
    total = len(prs.slides)
    
    results = {
        'file': pptx_path,
        'slides': total,
        'grade': grade,
        'template_expected': config['template'],
        'checks': {},
        'warnings': [],
        'pass': True,
    }
    
    if total == 0:
        results['error'] = 'No slides found'
        results['pass'] = False
        return results
    
    # ── Slide 1: Bloom objectives ──
    s1_texts = extract_text_from_slide(prs.slides[0])
    bloom_found = has_bloom_levels(s1_texts)
    results['checks']['bloom_objectives'] = {
        'passed': len(bloom_found) >= 2,
        'detail': f'Found: {bloom_found}' if bloom_found else 'Not detected',
    }
    if not results['checks']['bloom_objectives']['passed']:
        results['warnings'].append('Slide 1: Add Bloom learning objectives')
    
    # ── Words per slide ──
    word_counts = []
    over_limit = []
    for i, slide in enumerate(prs.slides):
        texts = extract_text_from_slide(slide)
        wc = count_words(texts)
        word_counts.append(wc)
        if wc > config['max_words']:
            over_limit.append((i + 1, wc))
    
    results['checks']['word_limit'] = {
        'passed': len(over_limit) == 0,
        'detail': f'Max: {config["max_words"]}/slide. Over: {len(over_limit)} slides' if over_limit else f'All under {config["max_words"]} words',
    }
    if over_limit:
        for num, wc in over_limit[:3]:
            results['warnings'].append(f'Slide {num}: {wc} words (limit {config["max_words"]})')
    
    # ── Paragraphs (no blocks >3 lines) ──
    long_blocks = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if len(para.text.split()) > 30 and '\n' not in para.text:
                        long_blocks.append((i + 1, len(para.text.split())))
    
    results['checks']['no_long_paragraphs'] = {
        'passed': len(long_blocks) == 0,
        'detail': f'{len(long_blocks)} long blocks found' if long_blocks else 'All text is chunked',
    }
    
    # ── Jargon detection ──
    all_texts = []
    for slide in prs.slides:
        all_texts.extend(extract_text_from_slide(slide))
    jargon_found = has_jargon(all_texts)
    
    results['checks']['jargon_handled'] = {
        'passed': len(jargon_found) == 0,
        'detail': f'Jargon terms: {jargon_found}' if jargon_found else 'No uncontrolled jargon',
    }
    if jargon_found:
        results['warnings'].append(f'Define in glossary: {", ".join(jargon_found[:4])}')
    
    # ── Questions ──
    total_questions = sum(has_questions(extract_text_from_slide(s)) for s in prs.slides)
    results['checks']['reflection_questions'] = {
        'passed': total_questions >= config['min_questions'],
        'detail': f'{total_questions} reflection triggers (min {config["min_questions"]})',
    }
    if not results['checks']['reflection_questions']['passed']:
        results['warnings'].append(f'Add ≥{config["min_questions"]} reflection questions')
    
    # ── Footer ──
    footer_found = any(has_footer(extract_text_from_slide(s)) for s in prs.slides)
    results['checks']['footer'] = {
        'passed': footer_found,
        'detail': 'Footer with grade+subject found' if footer_found else 'Not detected',
    }
    
    # ── Glossary (last slide) ──
    last_texts = extract_text_from_slide(prs.slides[-1])
    results['checks']['glossary'] = {
        'passed': has_glossary(last_texts),
        'detail': 'Glossary detected' if has_glossary(last_texts) else 'No glossary on last slide',
    }
    if not results['checks']['glossary']['passed']:
        results['warnings'].append('Add glossary to last slide')
    
    # ── Attribution ──
    attribution_found = has_attribution(all_texts)
    results['checks']['image_attribution'] = {
        'passed': attribution_found,
        'detail': 'CC attribution found' if attribution_found else 'No attribution detected',
    }
    
    # ── Template palette ──
    detected_templates = set()
    for slide in prs.slides:
        p = closest_palette(slide)
        if p:
            detected_templates.add(p)
    
    expected = config['template']
    results['checks']['template_palette'] = {
        'passed': expected in detected_templates,
        'detail': f'Expected: {expected}. Detected: {detected_templates or "none"}',
    }
    if not results['checks']['template_palette']['passed'] and detected_templates:
        results['warnings'].append(f'Palette {detected_templates} ≠ expected {expected}')
    
    # ── Fade transitions (progressive disclosure) ──
    fade_count = sum(1 for s in prs.slides if has_fade(s))
    results['checks']['fade_transitions'] = {
        'passed': fade_count >= 1,
        'detail': f'{fade_count} slides with fade transition',
    }
    
    # ── Overall ──
    failed = [k for k, v in results['checks'].items() if not v['passed']]
    results['pass'] = len(failed) == 0
    results['failed_checks'] = failed
    
    return results


# ── CLI ──────────────────────────────────────────────────────────────────────
def print_report(results):
    if 'error' in results:
        print(f'\n❌ ERROR: {results["error"]}')
        return
    
    status = 'PASS' if results['pass'] else 'FAIL'
    sep = '=' * 58
    print(f'\n{sep}')
    print(f'  VALIDATION: {status}')
    print(f'  File: {os.path.basename(results["file"])}')
    print(f'  Grade: {results["grade"]} | Slides: {results["slides"]}')
    print(f'  Expected template: {results["template_expected"]}')
    print(sep)
    print(f'\n  CHECK                         STATUS  DETAIL')
    print(f'  {"-"*55}')
    for check, data in results['checks'].items():
        icon = '[OK]' if data['passed'] else '[!!]'
        detail = data['detail'][:45]
        print(f'  {icon} {check:27s} {"PASS" if data["passed"] else "FAIL":10s} {detail}')
    
    if results['warnings']:
        print(f'\n  WARNINGS ({len(results["warnings"])}):')
        for w in results['warnings']:
            print(f'     - {w}')
    
    print(f'\n  Verdict: {status}')
    print(f'{sep}\n')
    return results['pass']


if __name__ == '__main__':
    args = sys.argv[1:]
    if not args or args[0] in ('-h', '--help'):
        print(__doc__)
        sys.exit(0)
    
    path = args[0]
    grade = '5to'
    if '--grade' in args:
        idx = args.index('--grade')
        if idx + 1 < len(args):
            grade = args[idx + 1]
    
    if not os.path.exists(path):
        print(f'❌ File not found: {path}')
        sys.exit(1)
    
    result = validate(path, grade)
    passed = print_report(result)
    sys.exit(0 if passed else 1)
